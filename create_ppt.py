"""
create_ppt.py — Generate a polished PowerPoint presentation
for the Cold-Start Energy Prediction project video presentation.

Strict grid system:
  - Slide: 13.333" × 7.5"  (widescreen 16:9)
  - Margins: left=0.7", right=0.7"  → content width = 11.933"
  - Top content starts at 1.55" (below header)
  - Gutter between columns: 0.35"
  - Footer zone: y >= 6.85"

Usage:
    python create_ppt.py
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ══════════════════════════════════════════
#  GRID CONSTANTS
# ══════════════════════════════════════════
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
MARGIN_L = 0.7
MARGIN_R = 0.7
CONTENT_W = SLIDE_W_IN - MARGIN_L - MARGIN_R   # 11.933
HEADER_Y = 0.35          # Title text top
ACCENT_BAR_Y = 1.05      # Underline accent bar
SUBTITLE_Y = 1.15        # Subtitle text
CONTENT_TOP = 1.55       # Where main content starts
FOOTER_Y = 6.95          # Footer zone top
GUTTER = 0.35            # Gap between columns

# Column widths (for 2-col and 3-col layouts)
COL2_W = (CONTENT_W - GUTTER) / 2              # ~5.79
COL3_W = (CONTENT_W - GUTTER * 2) / 3          # ~3.74
COL2_X = [MARGIN_L, MARGIN_L + COL2_W + GUTTER]
COL3_X = [MARGIN_L, MARGIN_L + COL3_W + GUTTER,
          MARGIN_L + (COL3_W + GUTTER) * 2]

# ══════════════════════════════════════════
#  COLOR PALETTE
# ══════════════════════════════════════════
DARK_BG       = RGBColor(0x0F, 0x17, 0x2A)
ACCENT_BLUE   = RGBColor(0x00, 0x96, 0xD6)
ACCENT_TEAL   = RGBColor(0x00, 0xB4, 0xD8)
ACCENT_GREEN  = RGBColor(0x2E, 0xCC, 0x71)
ACCENT_ORANGE = RGBColor(0xE6, 0x7E, 0x22)
ACCENT_RED    = RGBColor(0xE7, 0x4C, 0x3C)
ACCENT_PURPLE = RGBColor(0x9B, 0x59, 0xB6)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY    = RGBColor(0xCC, 0xCC, 0xCC)
MED_GRAY      = RGBColor(0x99, 0x99, 0x99)
CARD_BG       = RGBColor(0x16, 0x21, 0x3E)
HIGHLIGHT_BG  = RGBColor(0x1E, 0x3A, 0x5F)
GOLD          = RGBColor(0xFF, 0xD7, 0x00)
BORDER_DIM    = RGBColor(0x2A, 0x3A, 0x5C)
TRANSPARENT   = None

BASE_DIR    = os.path.dirname(os.path.abspath(__file__))
RESULTS_DIR = os.path.join(BASE_DIR, "results")
OUTPUT_PATH = os.path.join(BASE_DIR, "Cold_Start_Energy_Prediction_Final.pptx")

prs = Presentation()
prs.slide_width  = Inches(SLIDE_W_IN)
prs.slide_height = Inches(SLIDE_H_IN)


# ══════════════════════════════════════════
#  HELPER FUNCTIONS
# ══════════════════════════════════════════

def _no_border(shape):
    """Remove border from a shape."""
    shape.line.fill.background()


def _set_border(shape, color, width_pt=1.0):
    """Set a solid border on a shape."""
    shape.line.color.rgb = color
    shape.line.width = Pt(width_pt)


def dark_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG


def rect(slide, x, y, w, h, fill_color, border_color=None, border_pt=0):
    """Add a rectangle; returns the shape."""
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                               Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    if border_color:
        _set_border(s, border_color, border_pt)
    else:
        _no_border(s)
    return s


def rounded(slide, x, y, w, h, fill_color, border_color=None, border_pt=1.0):
    """Add a rounded rectangle."""
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    if border_color:
        _set_border(s, border_color, border_pt)
    else:
        _no_border(s)
    return s


def text(slide, x, y, w, h, txt, sz=18, color=WHITE, bold=False,
         align=PP_ALIGN.LEFT, font="Calibri", anchor=MSO_ANCHOR.TOP):
    """Add a text box with one paragraph. Returns the textbox shape."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.paragraphs[0].alignment = align
    except:
        pass
    p = tf.paragraphs[0]
    p.text = txt
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return tb


def bullets(slide, x, y, w, h, items, sz=14, color=LIGHT_GRAY,
            font="Calibri", spacing_pt=5):
    """Add multi-line bullet text. Each item is one paragraph."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(sz)
        p.font.color.rgb = color
        p.font.name = font
        p.space_after = Pt(spacing_pt)
    return tb


def img(slide, path, x, y, w=None, h=None):
    """Add an image if it exists."""
    if not os.path.exists(path):
        return None
    kw = {"image_file": path, "left": Inches(x), "top": Inches(y)}
    if w: kw["width"] = Inches(w)
    if h: kw["height"] = Inches(h)
    return slide.shapes.add_picture(**kw)


def slide_header(slide, title, subtitle=None, num=None):
    """Standard slide header: top accent bar, title, underline, subtitle, page#."""
    # Top accent bar — full width
    rect(slide, 0, 0, SLIDE_W_IN, 0.055, ACCENT_BLUE)
    # Title
    text(slide, MARGIN_L, HEADER_Y, CONTENT_W, 0.6, title,
         sz=30, color=WHITE, bold=True)
    # Underline
    rect(slide, MARGIN_L, ACCENT_BAR_Y, 2.2, 0.035, ACCENT_BLUE)
    # Subtitle
    if subtitle:
        text(slide, MARGIN_L, SUBTITLE_Y, CONTENT_W, 0.35, subtitle,
             sz=14, color=MED_GRAY)
    # Page number
    if num:
        text(slide, SLIDE_W_IN - 1.0, 7.05, 0.6, 0.3, str(num),
             sz=10, color=MED_GRAY, align=PP_ALIGN.RIGHT)


def card(slide, x, y, w, h, title_txt, items, accent=ACCENT_BLUE,
         tsz=16, isz=13, border=True):
    """Card container: rounded rect + left accent bar + title + bullets."""
    rounded(slide, x, y, w, h, CARD_BG,
            border_color=BORDER_DIM if border else None, border_pt=0.75)
    # Left accent bar
    rect(slide, x + 0.12, y + 0.12, 0.04, h - 0.24, accent)
    # Title
    text(slide, x + 0.28, y + 0.10, w - 0.40, 0.32, title_txt,
         sz=tsz, color=accent, bold=True)
    # Bullet items
    if items:
        bullets(slide, x + 0.28, y + 0.44, w - 0.40, h - 0.55,
                items, sz=isz, color=LIGHT_GRAY)


def badge(slide, x, y, w, h, txt, bg_color, txt_color=WHITE, sz=13):
    """Small rounded badge with centred text."""
    s = rounded(slide, x, y, w, h, bg_color)
    _no_border(s)
    tf = s.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = txt
    p.font.size = Pt(sz)
    p.font.color.rgb = txt_color
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    return s


# ══════════════════════════════════════════
#  SLIDE 1 — TITLE
# ══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)
rect(sl, 0, 0, SLIDE_W_IN, 0.07, ACCENT_BLUE)
rect(sl, 0, SLIDE_H_IN - 0.07, SLIDE_W_IN, 0.07, ACCENT_TEAL)
# Left accent line
rect(sl, 0.55, 1.6, 0.05, 4.2, ACCENT_BLUE)

text(sl, 0.9, 1.8, 11.5, 0.7,
     "Cold-Start Energy Consumption",
     sz=44, color=WHITE, bold=True)
text(sl, 0.9, 2.55, 11.5, 0.7,
     "Prediction for Smart Buildings",
     sz=44, color=ACCENT_BLUE, bold=True)
text(sl, 0.9, 3.55, 11, 0.45,
     "Using Machine Learning, Deep Learning & Transfer Learning Approaches",
     sz=20, color=LIGHT_GRAY)
rect(sl, 0.9, 4.25, 3.6, 0.03, ACCENT_TEAL)
text(sl, 0.9, 4.55, 10, 0.35,
     "Presented by: Neeraj Singh Bisht", sz=18, color=WHITE)
text(sl, 0.9, 5.0, 10, 0.30,
     "Dataset: CEEW Smart Meter Data — Mathura, India  |  2019 · 2020 · 2021",
     sz=14, color=MED_GRAY)
text(sl, 0.9, 5.35, 10, 0.30,
     "April 2026", sz=14, color=MED_GRAY)

# ══════════════════════════════════════════
#  SLIDE 2 — AGENDA
# ══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)
slide_header(sl, "Presentation Agenda", num=2)

agenda = [
    ("01", "Problem Statement & Motivation",            ACCENT_BLUE),
    ("02", "Dataset Description",                       ACCENT_TEAL),
    ("03", "Methodology & Pipeline Architecture",       ACCENT_GREEN),
    ("04", "Feature Engineering",                       ACCENT_ORANGE),
    ("05", "Models Implemented (ML · DL · Statistical)",ACCENT_PURPLE),
    ("06", "Transfer Learning — LSTM Approach",         ACCENT_BLUE),
    ("07", "Results & Performance Comparison",          ACCENT_TEAL),
    ("08", "Best Model Deep-Dive (HGBoost)",            GOLD),
    ("09", "Key Findings & Conclusions",                ACCENT_GREEN),
    ("10", "Future Work & References",                  ACCENT_ORANGE),
]

for i, (num, label, clr) in enumerate(agenda):
    yy = CONTENT_TOP + i * 0.52
    # Row background stripe
    rounded(sl, MARGIN_L, yy, CONTENT_W, 0.44, CARD_BG, border_color=BORDER_DIM, border_pt=0.5)
    # Number pill
    badge(sl, MARGIN_L + 0.12, yy + 0.05, 0.50, 0.34, num, clr, sz=13)
    # Label
    text(sl, MARGIN_L + 0.80, yy + 0.07, 10.5, 0.32, label,
         sz=17, color=WHITE)

# ══════════════════════════════════════════
#  SLIDE 3 — PROBLEM STATEMENT
# ══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)
slide_header(sl, "Problem Statement & Motivation",
             "Why cold-start prediction matters for smart buildings", num=3)

card(sl, COL2_X[0], CONTENT_TOP, COL2_W, 2.6,
     "THE COLD-START PROBLEM",
     ["Traditional forecasting requires historical data",
      "from the same building and time period.",
      "",
      "But this data doesn't exist for:",
      "  •  Newly installed smart meters",
      "  •  New buildings with no usage history",
      "  •  Post-renovation energy pattern changes"],
     accent=ACCENT_BLUE)

card(sl, COL2_X[1], CONTENT_TOP, COL2_W, 2.6,
     "OUR APPROACH",
     ["Train on 2019 + 2021 smart meter data",
      "Test on 2020 (completely unseen year)",
      "",
      "Simulates a real cold-start scenario",
      "where next-year data is unavailable",
      "",
      "Compare 12 models across 3 paradigms"],
     accent=ACCENT_GREEN)

# Research question highlight
rounded(sl, MARGIN_L, 4.45, CONTENT_W, 0.90, HIGHLIGHT_BG, border_color=ACCENT_BLUE, border_pt=1.0)
text(sl, MARGIN_L + 0.25, 4.55, CONTENT_W - 0.5, 0.7,
     "Research Question:  Can models trained on non-contiguous years (2019 + 2021) "
     "accurately predict energy consumption of an unseen year (2020)?",
     sz=16, color=GOLD, bold=True, align=PP_ALIGN.CENTER)

# Impact row
card(sl, MARGIN_L, 5.65, CONTENT_W, 1.15,
     "REAL-WORLD IMPACT",
     ["Enables energy planning for new buildings  ·  Supports smart-grid optimization  ·  "
      "Reduces energy waste  ·  Enables proactive demand response in smart cities"],
     accent=ACCENT_TEAL, tsz=15, isz=13)

# ══════════════════════════════════════════
#  SLIDE 4 — DATASET
# ══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)
slide_header(sl, "Dataset Description",
             "CEEW Smart Meter Data — Mathura, India", num=4)

ds_info = [
    ("2019 DATA", "CEEW Smart Meter\nMathura 2019",  "Training Set",       "~47 MB CSV",  ACCENT_BLUE),
    ("2020 DATA", "CEEW Smart Meter\nMathura 2020",  "Test Set (Unseen)",  "~182 MB CSV", ACCENT_RED),
    ("2021 DATA", "SM Cleaned Data\nMH 2021",        "Training Set",       "Supplementary", ACCENT_GREEN),
]

for i, (ttl, desc, role, sz_txt, clr) in enumerate(ds_info):
    x = COL3_X[i]
    rounded(sl, x, CONTENT_TOP, COL3_W, 2.55, CARD_BG, border_color=clr, border_pt=1.5)
    text(sl, x + 0.15, CONTENT_TOP + 0.12, COL3_W - 0.3, 0.30,
         ttl, sz=18, color=clr, bold=True, align=PP_ALIGN.CENTER)
    rect(sl, x + 0.4, CONTENT_TOP + 0.48, COL3_W - 0.8, 0.02, clr)  # mini divider
    text(sl, x + 0.15, CONTENT_TOP + 0.58, COL3_W - 0.3, 0.55,
         desc, sz=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    text(sl, x + 0.15, CONTENT_TOP + 1.25, COL3_W - 0.3, 0.30,
         f"Role: {role}", sz=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    text(sl, x + 0.15, CONTENT_TOP + 1.60, COL3_W - 0.3, 0.25,
         sz_txt, sz=11, color=MED_GRAY, align=PP_ALIGN.CENTER)
    # Bottom accent strip
    rect(sl, x, CONTENT_TOP + 2.50, COL3_W, 0.05, clr)

# Feature / split cards
card(sl, COL2_X[0], 4.45, COL2_W, 2.4,
     "KEY FEATURES",
     ["Target: t_kWh (total energy in kWh)",
      "Resolution: Hourly readings",
      "Preprocessing: Resampled to hourly mean",
      "Missing values: Linear interpolation",
      "Location: Mathura, Uttar Pradesh, India"],
     accent=ACCENT_BLUE, isz=13)

card(sl, COL2_X[1], 4.45, COL2_W, 2.4,
     "COLD-START SPLIT STRATEGY",
     ["Training: 2019 + 2021 combined data",
      "Testing: 2020 — completely unseen year",
      "No data leakage from the target year",
      "Simulates real-world deployment scenario",
      "Models must generalize across year gaps"],
     accent=ACCENT_ORANGE, isz=13)

# ══════════════════════════════════════════
#  SLIDE 5 — METHODOLOGY / ARCHITECTURE
# ══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)
slide_header(sl, "Methodology & Pipeline Architecture",
             "End-to-end prediction framework", num=5)

flow_path = os.path.join(BASE_DIR, "methodology_flow_diagram.jpg")
img(sl, flow_path, MARGIN_L, CONTENT_TOP, w=5.3, h=5.3)

steps = [
    ("1. Raw Data Collection",  "Load CEEW CSV files (2019, 2020, 2021)",                              ACCENT_BLUE),
    ("2. Data Preprocessing",   "Auto-detect time column · parse datetime · hourly resampling · interpolation", ACCENT_TEAL),
    ("3. Feature Engineering",  "Cyclical encoding (hour_sin/cos) · lag (1h, 24h) · sliding windows T=48",       ACCENT_GREEN),
    ("4. Cold-Start Split",     "Train: 2019 + 2021  →  Test: 2020 (unseen)",                          ACCENT_ORANGE),
    ("5. Model Training",       "12 models: 6 ML + 5 DL + 1 Statistical",                              ACCENT_PURPLE),
    ("6. Evaluation & Output",  "MAE · RMSE · R² · MAPE · CV(RMSE) + plots & CSV",                     ACCENT_RED),
]

STEP_H = 0.82
STEP_GAP = 0.06
RIGHT_X = 6.30
RIGHT_W = SLIDE_W_IN - MARGIN_R - RIGHT_X  # ≈6.33

for i, (sname, sdesc, sclr) in enumerate(steps):
    yy = CONTENT_TOP + i * (STEP_H + STEP_GAP)
    rounded(sl, RIGHT_X, yy, RIGHT_W, STEP_H, CARD_BG, border_color=BORDER_DIM, border_pt=0.5)
    # Accent strip on left edge of card
    rect(sl, RIGHT_X, yy, 0.05, STEP_H, sclr)
    text(sl, RIGHT_X + 0.20, yy + 0.06, RIGHT_W - 0.35, 0.28,
         sname, sz=14, color=sclr, bold=True)
    text(sl, RIGHT_X + 0.20, yy + 0.36, RIGHT_W - 0.35, 0.40,
         sdesc, sz=11, color=LIGHT_GRAY)

# ══════════════════════════════════════════
#  SLIDE 6 — FEATURE ENGINEERING
# ══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)
slide_header(sl, "Feature Engineering",
             "Transforming raw data into predictive features", num=6)

feat_data = [
    ("Cyclical Time Encoding", ACCENT_BLUE,
     ["hour_sin = sin(2π × hour / 24)",
      "hour_cos = cos(2π × hour / 24)",
      "",
      "Captures circular nature of time",
      "(Hour 23 is close to Hour 0)",
      "Prevents misleading linear features"]),
    ("Lag Features", ACCENT_GREEN,
     ["lag_1 : Energy value 1 hour ago",
      "lag_24: Energy value 24 hours ago",
      "",
      "Captures short-term patterns",
      "and daily periodicity in usage",
      "Recent history = best predictor"]),
    ("Sliding Window Sequences", ACCENT_ORANGE,
     ["Window size T = 48 hours (2 days)",
      "Used for LSTM, CNN, GRU, BiLSTM",
      "",
      "Input: 48 consecutive timesteps",
      "Output: Next timestep prediction",
      "Captures multi-day patterns"]),
]

for i, (ftitle, fclr, fitems) in enumerate(feat_data):
    card(sl, COL3_X[i], CONTENT_TOP, COL3_W, 4.1,
         ftitle, fitems, accent=fclr, tsz=15, isz=13)

# Insight bar
rounded(sl, MARGIN_L, 5.9, CONTENT_W, 0.75, HIGHLIGHT_BG, border_color=GOLD, border_pt=0.75)
text(sl, MARGIN_L + 0.2, 5.98, CONTENT_W - 0.4, 0.6,
     "Key Insight:  Combining cyclical encoding with lag features lets models learn "
     "time-of-day patterns and recent consumption trends — critical for cold-start generalization.",
     sz=13, color=GOLD, bold=True, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════
#  SLIDE 7 — MODELS OVERVIEW
# ══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)
slide_header(sl, "Models Implemented — Overview",
             "12 models across 3 paradigms", num=7)

cats = [
    ("MACHINE LEARNING (6)", ACCENT_BLUE,
     ["Linear Regression", "Ridge Regression", "Lasso Regression",
      "XGBoost (Gradient Boosting)", "SVR (Support Vector)", "HGBoost (AutoML)"]),
    ("DEEP LEARNING (5)", ACCENT_PURPLE,
     ["LSTM (Long Short-Term Memory)", "CNN (1D Convolutional)",
      "GRU (Gated Recurrent Unit)", "BiLSTM (Bidirectional)",
      "Transfer LSTM (Fine-tune)"]),
    ("STATISTICAL (1)", ACCENT_ORANGE,
     ["ARIMA (AutoRegressive", "  Integrated Moving Avg)"]),
]

for ci, (cat_title, cat_clr, models) in enumerate(cats):
    x = COL3_X[ci]
    # Header bar
    rounded(sl, x, CONTENT_TOP, COL3_W, 0.48, cat_clr)
    text(sl, x + 0.1, CONTENT_TOP + 0.06, COL3_W - 0.2, 0.36,
         cat_title, sz=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # Model list items
    for mi, mname in enumerate(models):
        yy = CONTENT_TOP + 0.62 + mi * 0.50
        rounded(sl, x, yy, COL3_W, 0.43, CARD_BG, border_color=BORDER_DIM, border_pt=0.5)
        text(sl, x + 0.15, yy + 0.06, COL3_W - 0.3, 0.30,
             f"▸  {mname}", sz=12, color=LIGHT_GRAY)

# Footer summary
text(sl, MARGIN_L, FOOTER_Y, CONTENT_W, 0.35,
     "Total: 6 ML  +  5 DL  +  1 Statistical  =  12 models compared under cold-start conditions",
     sz=15, color=ACCENT_TEAL, bold=True, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════
#  SLIDE 8 — ML MODELS DETAIL
# ══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)
slide_header(sl, "Machine Learning Models — Details",
             "Traditional and ensemble approaches", num=8)

ml_data = [
    ("Linear Regression", "Baseline OLS model",           "R² = 0.9902", ACCENT_BLUE),
    ("Ridge Regression",  "L2 regularization, α = 1.0",   "R² = 0.9900", ACCENT_TEAL),
    ("Lasso Regression",  "L1 regularization, α = 0.01",  "R² = 0.3105", ACCENT_ORANGE),
    ("XGBoost",           "600 trees, lr=0.03, depth=4",   "R² = 0.9912", ACCENT_GREEN),
    ("SVR",               "RBF kernel, C=10, ε=0.1",      "R² = 0.9912", ACCENT_PURPLE),
    ("HGBoost",           "Bayesian optimized, 250 evals", "R² = 0.9922", GOLD),
]

for i, (mn, md, mr, mc) in enumerate(ml_data):
    row, col = divmod(i, 3)
    x = COL3_X[col]
    yy = CONTENT_TOP + row * 2.70

    rounded(sl, x, yy, COL3_W, 2.30, CARD_BG, border_color=mc, border_pt=1.2)
    text(sl, x + 0.18, yy + 0.14, COL3_W - 0.36, 0.30,
         mn, sz=17, color=mc, bold=True)
    rect(sl, x + 0.18, yy + 0.50, COL3_W * 0.6, 0.02, mc)
    text(sl, x + 0.18, yy + 0.60, COL3_W - 0.36, 0.65,
         md, sz=12, color=LIGHT_GRAY)
    badge(sl, x + 0.18, yy + 1.50, 1.70, 0.42, mr, mc, sz=13)

# ══════════════════════════════════════════
#  SLIDE 9 — DL MODELS DETAIL
# ══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)
slide_header(sl, "Deep Learning Models — Details",
             "Sequence-based neural network architectures", num=9)

dl_data = [
    ("LSTM",          "64 units, 32 dense, 80 epochs\nCaptures long-term dependencies",  "R² = 0.9146", ACCENT_PURPLE),
    ("CNN (1D)",      "64 filters, kernel=3, pool=2\n1D convolutions for local patterns", "R² = 0.8673", ACCENT_ORANGE),
    ("GRU",           "64 units, 80 epochs\nLighter LSTM alternative",                    "R² = 0.8815", ACCENT_TEAL),
    ("BiLSTM",        "64 units bidirectional\nForward + backward context",               "R² = 0.8703", ACCENT_RED),
    ("Transfer LSTM", "Pre-train 2019 → Fine-tune 2021\n3-phase progressive unfreezing", "R² = 0.8866", ACCENT_BLUE),
]

# Row 1: first 3
for i in range(3):
    mn, md, mr, mc = dl_data[i]
    x = COL3_X[i]
    rounded(sl, x, CONTENT_TOP, COL3_W, 2.30, CARD_BG, border_color=mc, border_pt=1.2)
    text(sl, x + 0.18, CONTENT_TOP + 0.14, COL3_W - 0.36, 0.30,
         mn, sz=17, color=mc, bold=True)
    rect(sl, x + 0.18, CONTENT_TOP + 0.50, COL3_W * 0.6, 0.02, mc)
    text(sl, x + 0.18, CONTENT_TOP + 0.60, COL3_W - 0.36, 0.75,
         md, sz=12, color=LIGHT_GRAY)
    badge(sl, x + 0.18, CONTENT_TOP + 1.55, 1.70, 0.42, mr, mc, sz=13)

# Row 2: last 2 centred
row2_y = CONTENT_TOP + 2.55
ROW2_X = [MARGIN_L + (CONTENT_W - 2 * COL3_W - GUTTER) / 2,
          MARGIN_L + (CONTENT_W - 2 * COL3_W - GUTTER) / 2 + COL3_W + GUTTER]

for j in range(2):
    mn, md, mr, mc = dl_data[3 + j]
    x = ROW2_X[j]
    rounded(sl, x, row2_y, COL3_W, 2.30, CARD_BG, border_color=mc, border_pt=1.2)
    text(sl, x + 0.18, row2_y + 0.14, COL3_W - 0.36, 0.30,
         mn, sz=17, color=mc, bold=True)
    rect(sl, x + 0.18, row2_y + 0.50, COL3_W * 0.6, 0.02, mc)
    text(sl, x + 0.18, row2_y + 0.60, COL3_W - 0.36, 0.75,
         md, sz=12, color=LIGHT_GRAY)
    badge(sl, x + 0.18, row2_y + 1.55, 1.70, 0.42, mr, mc, sz=13)

# Footer config
text(sl, MARGIN_L, FOOTER_Y, CONTENT_W, 0.30,
     "Common config:  Sequence length = 48 h  |  Optimizer = Adam  |  Loss = MSE  |  Batch = 64",
     sz=12, color=MED_GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════
#  SLIDE 10 — TRANSFER LEARNING
# ══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)
slide_header(sl, "Transfer Learning — LSTM Approach",
             "Progressive unfreezing for cold-start adaptation", num=10)

phases = [
    ("PHASE 1", "Pre-training on Source", ACCENT_BLUE,
     ["Train LSTM on 2019 data",
      "64 units, 32 dense",
      "50 epochs, Adam optimizer",
      "Learns general energy patterns"]),
    ("PHASE 2", "Fine-tune Dense Head", ACCENT_ORANGE,
     ["Freeze LSTM layers",
      "Fine-tune output dense only",
      "Train on 2021 (target domain)",
      "20 epochs, LR = 5×10⁻⁴"]),
    ("PHASE 3", "Full Fine-tuning", ACCENT_GREEN,
     ["Unfreeze all layers",
      "Very low LR: 1×10⁻⁵",
      "20 epochs gentle tuning",
      "Prevents catastrophic forgetting"]),
]

for i, (ph_num, ph_title, ph_clr, ph_items) in enumerate(phases):
    x = COL3_X[i]
    # Phase label bar
    rounded(sl, x, CONTENT_TOP, COL3_W, 0.55, ph_clr)
    text(sl, x + 0.1, CONTENT_TOP + 0.02, COL3_W - 0.2, 0.22,
         ph_num, sz=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    text(sl, x + 0.1, CONTENT_TOP + 0.24, COL3_W - 0.2, 0.28,
         ph_title, sz=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # Arrow between phases
    if i < 2:
        text(sl, x + COL3_W + 0.02, CONTENT_TOP + 0.10, GUTTER - 0.04, 0.40,
             "→", sz=24, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    # Detail card
    card(sl, x, CONTENT_TOP + 0.70, COL3_W, 2.40,
         "", ph_items, accent=ph_clr, tsz=14, isz=13)

# Result highlight
rounded(sl, MARGIN_L, 4.95, CONTENT_W, 1.80, HIGHLIGHT_BG, border_color=ACCENT_BLUE, border_pt=1.0)
text(sl, MARGIN_L + 0.25, 5.05, CONTENT_W - 0.5, 0.30,
     "Transfer LSTM Results", sz=18, color=ACCENT_BLUE, bold=True)
text(sl, MARGIN_L + 0.25, 5.40, CONTENT_W - 0.5, 0.30,
     "R² = 0.8866   |   MAE = 0.0015   |   RMSE = 0.0028   |   MAPE = 15.80%",
     sz=16, color=WHITE, bold=True)
text(sl, MARGIN_L + 0.25, 5.80, CONTENT_W - 0.5, 0.75,
     "While not outperforming traditional ML models, Transfer LSTM demonstrates the viability "
     "of cross-domain knowledge transfer for cold-start scenarios. With more diverse source data "
     "and advanced domain-adaptation techniques, this gap can be further narrowed.",
     sz=13, color=LIGHT_GRAY)

# ══════════════════════════════════════════
#  SLIDE 11 — RESULTS TABLE
# ══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)
slide_header(sl, "Results — Full Performance Comparison",
             "All 12 models ranked by R² (descending)", num=11)

headers = ["Model", "MAE", "RMSE", "R²", "MAPE (%)", "CV(RMSE) (%)"]
rows = [
    ["HGBoost  ★",         "0.0004", "0.0007", "0.9922", "4.53",   "8.44"],
    ["XGBoost",            "0.0004", "0.0008", "0.9912", "4.67",   "9.00"],
    ["SVR",                "0.0005", "0.0008", "0.9912", "7.49",   "8.98"],
    ["Linear Regression",  "0.0006", "0.0008", "0.9902", "8.09",   "9.46"],
    ["Ridge Regression",   "0.0006", "0.0008", "0.9900", "8.50",   "9.60"],
    ["LSTM",               "0.0012", "0.0025", "0.9146", "13.10", "28.04"],
    ["Transfer LSTM",      "0.0015", "0.0028", "0.8866", "15.80", "32.31"],
    ["GRU",                "0.0013", "0.0029", "0.8815", "13.09", "33.03"],
    ["BiLSTM",             "0.0015", "0.0030", "0.8703", "15.09", "34.55"],
    ["CNN (1D)",           "0.0016", "0.0031", "0.8673", "17.59", "34.96"],
    ["Lasso Regression",   "0.0063", "0.0069", "0.3105", "115.16","79.51"],
    ["ARIMA",              "0.0050", "0.0089", "−0.14",  "44.26", "102.30"],
]

tbl_shape = sl.shapes.add_table(
    len(rows) + 1, len(headers),
    Inches(MARGIN_L), Inches(CONTENT_TOP),
    Inches(CONTENT_W), Inches(5.0))
tbl = tbl_shape.table

# Column widths  (approximate relative)
col_widths = [2.8, 1.5, 1.5, 1.5, 1.8, 2.0]  # must sum ≈ CONTENT_W but pptx ignores exact
total = sum(col_widths)
for ci, cw in enumerate(col_widths):
    tbl.columns[ci].width = Emu(int(Inches(CONTENT_W).emu * cw / total))

# Header row
for ci, h in enumerate(headers):
    cell = tbl.cell(0, ci)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = ACCENT_BLUE
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(11)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER

# Data rows
TOP_BG     = [RGBColor(0x14, 0x2E, 0x22), RGBColor(0x14, 0x2B, 0x20), RGBColor(0x14, 0x28, 0x1E)]
BOTTOM_BG  = RGBColor(0x30, 0x18, 0x18)

for ri, row in enumerate(rows):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri + 1, ci)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            p.font.bold = (ci == 0)
            if ri < 3:
                p.font.color.rgb = ACCENT_GREEN
            elif ri < 5:
                p.font.color.rgb = LIGHT_GRAY
            elif ri < 10:
                p.font.color.rgb = MED_GRAY
            else:
                p.font.color.rgb = ACCENT_RED
        # Background
        cell.fill.solid()
        if ri < 3:
            cell.fill.fore_color.rgb = TOP_BG[ri]
        elif ri >= 10:
            cell.fill.fore_color.rgb = BOTTOM_BG
        else:
            cell.fill.fore_color.rgb = CARD_BG

text(sl, MARGIN_L, FOOTER_Y, CONTENT_W, 0.30,
     "★ Best overall   |   Green = Top performers   |   Red = Poor performers",
     sz=11, color=MED_GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════
#  SLIDE 12 — VISUAL: MODEL COMPARISON
# ══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)
slide_header(sl, "Visual Results — Prediction Comparison",
             "ML models (top)  vs  Deep Learning models (bottom)", num=12)

img(sl, os.path.join(RESULTS_DIR, "model_comparison.png"),
    0.2, 1.40, w=12.9, h=5.9)

# ══════════════════════════════════════════
#  SLIDE 13 — VISUAL: METRICS BAR
# ══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)
slide_header(sl, "Visual Results — Metrics Comparison",
             "MAE · RMSE · R²  across all 12 models", num=13)

img(sl, os.path.join(RESULTS_DIR, "metrics_comparison.png"),
    0.2, 1.50, w=12.9, h=4.3)

rounded(sl, MARGIN_L, 6.1, CONTENT_W, 0.70, HIGHLIGHT_BG, border_color=GOLD, border_pt=0.75)
text(sl, MARGIN_L + 0.2, 6.18, CONTENT_W - 0.4, 0.55,
     "ML models (HGBoost, XGBoost, SVR) significantly outperform DL models in the cold-start "
     "scenario. ARIMA and Lasso show poor generalization to unseen years.",
     sz=13, color=GOLD, bold=True, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════
#  SLIDE 14 — BEST MODEL: HGBoost Scatter
# ══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)
slide_header(sl, "Best Model — HGBoost (Actual vs Predicted)",
             "R² = 0.9922  |  MAE = 0.0004  |  RMSE = 0.0007  |  MAPE = 4.53%", num=14)

img(sl, os.path.join(RESULTS_DIR, "actual_vs_predicted_HGBoost.png"),
    0.3, 1.45, w=12.7, h=5.2)

text(sl, MARGIN_L, FOOTER_Y, CONTENT_W, 0.30,
     "HGBoost achieves the best R² and lowest error across all metrics under cold-start conditions.",
     sz=13, color=GOLD, bold=True, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════
#  SLIDE 15 — BEST ML: XGBoost Scatter
# ══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)
slide_header(sl, "Runner-up — XGBoost (Actual vs Predicted)",
             "R² = 0.9912  |  MAE = 0.0004  |  RMSE = 0.0008  |  MAPE = 4.67%", num=15)

img(sl, os.path.join(RESULTS_DIR, "actual_vs_predicted_XGBoost.png"),
    0.3, 1.45, w=12.7, h=5.2)

text(sl, MARGIN_L, FOOTER_Y, CONTENT_W, 0.30,
     "XGBoost matches HGBoost closely with gradient boosting and is the fastest high-accuracy model.",
     sz=13, color=ACCENT_GREEN, bold=True, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════
#  SLIDE 16 — BEST DL: LSTM Scatter
# ══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)
slide_header(sl, "Best DL Model — LSTM (Actual vs Predicted)",
             "R² = 0.9146  |  MAE = 0.0012  |  RMSE = 0.0025  |  MAPE = 13.10%", num=16)

img(sl, os.path.join(RESULTS_DIR, "actual_vs_predicted_LSTM.png"),
    0.3, 1.45, w=12.7, h=5.2)

text(sl, MARGIN_L, FOOTER_Y, CONTENT_W, 0.30,
     "LSTM is the best deep learning model, though it trails ML models — DL needs more cold-start data.",
     sz=13, color=ACCENT_PURPLE, bold=True, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════
#  SLIDE 17 — LR FAMILY COMPARISON
# ══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)
slide_header(sl, "Linear Regression Family — Comparison",
             "Linear vs Ridge vs Lasso under cold-start conditions", num=17)

img(sl, os.path.join(RESULTS_DIR, "lr_ridge_lasso_comparison.png"),
    0.3, 1.45, w=12.7, h=5.2)

text(sl, MARGIN_L, FOOTER_Y, CONTENT_W, 0.30,
     "Linear & Ridge perform comparably (R² ≈ 0.99). Lasso's aggressive L1 regularization destroys features needed for cold-start.",
     sz=12, color=ACCENT_ORANGE, bold=True, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════
#  SLIDE 18 — KEY FINDINGS
# ══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)
slide_header(sl, "Key Findings & Conclusions", num=18)

findings = [
    ("HGBoost is the Best Overall Model",
     "R² = 0.9922, lowest CV(RMSE) of 8.44%. Bayesian hyperparameter optimization delivers the best cold-start performance.",
     GOLD),
    ("ML Models Dominate Cold-Start Prediction",
     "Top 5 models are all ML-based (R² > 0.99). Tabular feature-engineered models outperform sequence-based DL in this scenario.",
     ACCENT_GREEN),
    ("Deep Learning Achieves R² ≈ 0.87–0.91",
     "LSTM leads DL (0.9146). DL requires more diverse training data to match ML in cold-start conditions.",
     ACCENT_PURPLE),
    ("Transfer Learning Shows Promise",
     "Transfer LSTM (R²=0.8866) validates cross-domain knowledge transfer despite limited source diversity.",
     ACCENT_BLUE),
    ("ARIMA & Lasso Fail in Cold-Start",
     "ARIMA (R²=−0.14) and Lasso (R²=0.31) cannot generalize across years. Not suitable for cold-start deployment.",
     ACCENT_RED),
]

ROW_H = 0.98
for i, (ftitle, fdesc, fclr) in enumerate(findings):
    yy = CONTENT_TOP + i * (ROW_H + 0.08)
    rounded(sl, MARGIN_L, yy, CONTENT_W, ROW_H, CARD_BG, border_color=BORDER_DIM, border_pt=0.5)
    rect(sl, MARGIN_L, yy, 0.06, ROW_H, fclr)
    text(sl, MARGIN_L + 0.22, yy + 0.06, CONTENT_W - 0.40, 0.28,
         ftitle, sz=15, color=fclr, bold=True)
    text(sl, MARGIN_L + 0.22, yy + 0.38, CONTENT_W - 0.40, 0.55,
         fdesc, sz=12, color=LIGHT_GRAY)

# ══════════════════════════════════════════
#  SLIDE 19 — FUTURE WORK
# ══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)
slide_header(sl, "Future Work & Improvements", num=19)

future = [
    ("Ensemble Stacking",          "Combine HGBoost + LSTM to leverage both ML and DL strengths.",     ACCENT_BLUE),
    ("Weather Data Integration",   "Add temperature, humidity, solar radiation for richer features.",    ACCENT_GREEN),
    ("Transformer Attention",      "Replace RNN layers with self-attention for long-range patterns.",    ACCENT_PURPLE),
    ("Multi-Building Transfer",    "Extend transfer learning across buildings & geographic regions.",     ACCENT_ORANGE),
    ("Real-Time Deployment",       "Online learning pipeline with streaming predictions for smart grids.", ACCENT_TEAL),
    ("Explainability (XAI)",       "Apply SHAP / LIME for model interpretability in energy decisions.",  ACCENT_RED),
]

for i, (ft, fd, fc) in enumerate(future):
    row, col = divmod(i, 2)
    x = COL2_X[col]
    yy = CONTENT_TOP + row * 1.75
    rounded(sl, x, yy, COL2_W, 1.45, CARD_BG, border_color=fc, border_pt=1.0)
    rect(sl, x + 0.12, yy + 0.14, 0.04, 1.17, fc)
    text(sl, x + 0.30, yy + 0.12, COL2_W - 0.5, 0.30,
         f"→  {ft}", sz=15, color=fc, bold=True)
    text(sl, x + 0.30, yy + 0.50, COL2_W - 0.5, 0.85,
         fd, sz=12, color=LIGHT_GRAY)

# ══════════════════════════════════════════
#  SLIDE 20 — REFERENCES
# ══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)
slide_header(sl, "References & Tools", num=20)

refs = [
    "CEEW (Council on Energy, Environment and Water) — Smart Meter Data, Mathura, India",
    "Hochreiter & Schmidhuber (1997) — Long Short-Term Memory (LSTM)",
    "Chen & Guestrin (2016) — XGBoost: A Scalable Tree Boosting System",
    "Cho et al. (2014) — Learning Phrase Representations using RNN Encoder–Decoder (GRU)",
    "Yosinski et al. (2014) — How Transferable are Features in Deep Neural Networks?",
    "Box, Jenkins, Reinsel (2015) — Time Series Analysis: ARIMA Models",
]

tools = [
    "Python 3.8+  ·  TensorFlow / Keras  ·  scikit-learn  ·  XGBoost",
    "pandas  ·  NumPy  ·  Matplotlib  ·  pmdarima  ·  HGBoost",
    "GitHub Repository: github.com/Nsb05/Cold-Heart",
]

card(sl, MARGIN_L, CONTENT_TOP, CONTENT_W, 3.0,
     "REFERENCES", refs, accent=ACCENT_BLUE, tsz=17, isz=12)

card(sl, MARGIN_L, CONTENT_TOP + 3.25, CONTENT_W, 1.95,
     "TOOLS & TECHNOLOGIES", tools, accent=ACCENT_TEAL, tsz=17, isz=13)

text(sl, MARGIN_L, FOOTER_Y, CONTENT_W, 0.30,
     "All code is open-source under MIT License",
     sz=12, color=MED_GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════
#  SLIDE 21 — THANK YOU
# ══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(sl)
rect(sl, 0, 0, SLIDE_W_IN, 0.07, ACCENT_BLUE)
rect(sl, 0, SLIDE_H_IN - 0.07, SLIDE_W_IN, 0.07, ACCENT_TEAL)

text(sl, 0, 2.0, SLIDE_W_IN, 0.9,
     "Thank You!", sz=54, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
rect(sl, (SLIDE_W_IN - 3.0) / 2, 3.1, 3.0, 0.035, ACCENT_BLUE)
text(sl, 0, 3.4, SLIDE_W_IN, 0.5,
     "Cold-Start Energy Consumption Prediction for Smart Buildings",
     sz=19, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
text(sl, 0, 4.1, SLIDE_W_IN, 0.45,
     "Neeraj Singh Bisht", sz=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
text(sl, 0, 4.7, SLIDE_W_IN, 0.35,
     "github.com/Nsb05/Cold-Heart", sz=14, color=MED_GRAY, align=PP_ALIGN.CENTER)
text(sl, 0, 5.3, SLIDE_W_IN, 0.4,
     "Questions & Discussion", sz=18, color=ACCENT_TEAL, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════
#  SAVE
# ══════════════════════════════════════════
prs.save(OUTPUT_PATH)

slide_list = [
    "Title Slide",
    "Agenda",
    "Problem Statement & Motivation",
    "Dataset Description",
    "Methodology & Pipeline Architecture",
    "Feature Engineering",
    "Models Implemented — Overview",
    "ML Models — Details",
    "DL Models — Details",
    "Transfer Learning — LSTM Approach",
    "Results — Full Performance Table",
    "Visual: Prediction Comparison",
    "Visual: Metrics Bar Charts",
    "Best Model — HGBoost (Actual vs Predicted)",
    "Runner-up — XGBoost (Actual vs Predicted)",
    "Best DL — LSTM (Actual vs Predicted)",
    "LR Family Comparison",
    "Key Findings & Conclusions",
    "Future Work & Improvements",
    "References & Tools",
    "Thank You",
]

print(f"\n✅ Presentation saved to: {OUTPUT_PATH}")
print(f"   Total slides: {len(prs.slides)}\n")
for i, s in enumerate(slide_list, 1):
    print(f"   {i:2d}. {s}")
print()
